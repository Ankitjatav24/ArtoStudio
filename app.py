import streamlit as st
import pandas as pd
import io
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import plotly.express as px
import matplotlib.pyplot as plt
import seaborn as sns

# --- Third-Party Libraries Safe Import ---
try:
    from twilio.rest import Client
except:
    pass

try:
    from fpdf import FPDF
except:
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except:
    pass

# ==========================================
# 1. NEXT-GEN ANIMATED UI/UX CUSTOM CSS
# ==========================================
st.set_page_config(page_title="Arto+ Automation Studio Pro", page_icon="⚡", layout="wide")

st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght=300;400;500;600;700;800&display=swap');
    
    /* --- ANIMATION KEYFRAMES --- */
    @keyframes fadeInUp {
        0% { opacity: 0; transform: translateY(24px); filter: blur(4px); }
        100% { opacity: 1; transform: translateY(0); filter: blur(0); }
    }
    
    @keyframes pulseGlow {
        0% { box-shadow: 0 10px 25px -5px rgba(37, 99, 235, 0.05); }
        50% { box-shadow: 0 20px 35px 0 rgba(37, 99, 235, 0.15); }
        100% { box-shadow: 0 10px 25px -5px rgba(37, 99, 235, 0.05); }
    }
    
    @keyframes gradientShimmer {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    @keyframes goldGlow {
        0% { border-color: rgba(234, 179, 8, 0.3); box-shadow: 0 4px 12px rgba(234, 179, 8, 0.05); }
        50% { border-color: rgba(234, 179, 8, 1); box-shadow: 0 10px 25px rgba(234, 179, 8, 0.3); }
        100% { border-color: rgba(234, 179, 8, 0.3); box-shadow: 0 4px 12px rgba(234, 179, 8, 0.05); }
    }

    /* Global Container Tuning & Fluid Entry */
    html, body, [data-testid="stAppViewContainer"] {
        font-family: 'Plus Jakarta Sans', sans-serif !important;
        background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%) !important;
    }
    
    /* Apply Master Animation to the Main Block */
    [data-testid="stVerticalBlock"] > div {
        animation: fadeInUp 0.6s cubic-bezier(0.16, 1, 0.3, 1) both;
    }
    
    /* Header Typography Core */
    h1, h2, h3, h4 { 
        color: #0f172a !important; 
        font-weight: 800 !important;
        letter-spacing: -0.5px !important;
    }
    
    /* Animated Glassmorphic Cards for Metrics */
    [data-testid="stMetric"] {
        background: rgba(255, 255, 255, 0.7) !important;
        backdrop-filter: blur(12px) !important;
        -webkit-backdrop-filter: blur(12px) !important;
        padding: 24px !important;
        border-radius: 16px !important;
        border: 1px solid rgba(255, 255, 255, 0.6) !important;
        animation: pulseGlow 4s infinite ease-in-out;
        transition: transform 0.4s cubic-bezier(0.16, 1, 0.3, 1), box-shadow 0.4s cubic-bezier(0.16, 1, 0.3, 1) !important;
    }
    [data-testid="stMetric"]:hover {
        transform: translateY(-6px) scale(1.01) !important;
        box-shadow: 0 25px 30px -5px rgba(37, 99, 235, 0.12), 0 12px 15px -5px rgba(37, 99, 235, 0.06) !important;
    }
    [data-testid="stMetricValue"] { 
        color: #1e40af !important; 
        font-weight: 800 !important; 
        font-size: 32px !important;
    }
    
    /* Premium Cyber-Gradient Tabs with Active Underline Animation */
    button[data-baseweb="tab"] { 
        font-size: 14px !important; 
        font-weight: 700 !important; 
        padding: 12px 20px !important;
        color: #64748b !important;
        border: none !important;
        background-color: transparent !important;
        transition: all 0.3s cubic-bezier(0.16, 1, 0.3, 1) !important;
    }
    button[data-baseweb="tab"]:hover {
        color: #2563eb !important;
        transform: translateY(-1px);
    }
    button[data-baseweb="tab"][aria-selected="true"] { 
        background: #ffffff !important;
        color: #2563eb !important; 
        border-radius: 12px 12px 0 0 !important;
        box-shadow: 0 -4px 12px rgba(0, 0, 0, 0.02) !important;
        border-bottom: 3px solid #2563eb !important;
    }
    
    /* Kinetic Luxury Interactive Buttons */
    .stButton>button { 
        background: linear-gradient(135deg, #2563eb 0%, #1d4ed8 100%) !important; 
        color: white !important; 
        font-weight: 700 !important; 
        letter-spacing: 0.3px !important;
        border-radius: 12px !important; 
        border: none !important;
        padding: 12px 24px !important;
        box-shadow: 0 4px 12px rgba(37, 99, 235, 0.15) !important;
        transition: all 0.3s cubic-bezier(0.16, 1, 0.3, 1) !important;
        width: 100%;
    }
    .stButton>button:hover {
        background: linear-gradient(135deg, #1d4ed8 0%, #1e40af 100%) !important;
        box-shadow: 0 8px 24px rgba(37, 99, 235, 0.35) !important;
        transform: translateY(-3px) scale(1.02) !important;
    }
    .stButton>button:active {
        transform: translateY(0px) scale(0.98) !important;
    }
    
    /* Top Luxury Fluid Branding Navbar */
    .brand-top-bar {
        background: linear-gradient(90deg, #ffffff 0%, #f8fafc 100%); 
        border: 1px solid rgba(226, 232, 240, 0.8);
        padding: 24px 32px; 
        border-radius: 20px;
        margin-bottom: 30px;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.02);
        display: flex;
        justify-content: space-between;
        align-items: center;
        animation: fadeInUp 0.5s cubic-bezier(0.16, 1, 0.3, 1) both;
    }
    .brand-title { font-size: 30px; font-weight: 800; color: #2563eb; letter-spacing: -1px; }
    .brand-subtitle { color: #0f172a; font-weight: 500; font-size: 18px; margin-left: 8px;}
    
    /* Animated Pro Shimmer Badge */
    .badge-pro { 
        background: linear-gradient(-45deg, #3b82f6, #1d4ed8, #2563eb, #1e40af);
        background-size: 300% 300%;
        animation: gradientShimmer 4s ease infinite;
        color: white; 
        padding: 6px 14px; 
        border-radius: 99px; 
        font-size: 12px; 
        font-weight: 700; 
        text-transform: uppercase; 
        letter-spacing: 1px;
    }
    
    /* Left Sidebar Re-skinning */
    [data-testid="stSidebar"] {
        background-color: #ffffff !important;
        border-right: 1px solid #e2e8f0 !important;
    }
    
    /* Smooth Transitions for Interactive Data Editors */
    .stDataEditor {
        border-radius: 14px !important;
        overflow: hidden !important;
        box-shadow: 0 4px 12px rgba(0,0,0,0.01) !important;
        transition: all 0.3s ease !important;
    }
    .stDataEditor:hover {
        box-shadow: 0 12px 24px rgba(0,0,0,0.03) !important;
    }

    /* PREMIUM GATE LOCK VISUALS */
    .premium-lock-container {
        background: rgba(254, 243, 199, 0.4) !important;
        border: 2px dashed #eab308 !important;
        padding: 24px;
        border-radius: 16px;
        text-align: center;
        animation: goldGlow 3s infinite ease-in-out;
        margin-top: 15px;
    }
    .premium-lock-title {
        color: #854d0e !important;
        font-weight: 800 !important;
        font-size: 18px;
        margin-bottom: 6px;
    }
    .premium-lock-desc {
        color: #a16207 !important;
        font-size: 13px;
        margin-bottom: 0px;
    }
    .premium-blur-zone {
        filter: blur(5px);
        pointer-events: none;
        opacity: 0.4;
        user-select: none;
    }
    </style>
""", unsafe_allow_html=True)


# ==========================================
# 2. ALL-GRAPH SEAMLESS REPORT GENERATOR
# ==========================================
def generate_report_graph_bytes(df, x_col, y_col, chart_type):
    try:
        fig, ax = plt.subplots(figsize=(10, 6))
        sns.set_theme(style="whitegrid")
        
        sample_df = df.head(40).copy()
        sample_df[x_col] = sample_df[x_col].astype(str)
        
        if y_col:
            sample_df[y_col] = pd.to_numeric(sample_df[y_col], errors='coerce').fillna(0)

        if chart_type in ["Bar Diagram", "Bar"]:
            sns.barplot(data=sample_df, x=x_col, y=y_col, palette="Blues_d", ax=ax)
        elif chart_type in ["Line Plot Matrix", "Line"]:
            sns.lineplot(data=sample_df, x=x_col, y=y_col, marker="o", color="#2563eb", linewidth=2, ax=ax)
        elif chart_type in ["Scatter Cluster", "Scatter"]:
            sns.scatterplot(data=sample_df, x=x_col, y=y_col, color="#2563eb", s=100, ax=ax)
        elif chart_type in ["Area Horizon", "Area"]:
            ax.fill_between(sample_df[x_col], sample_df[y_col], color="#2563eb", alpha=0.4)
            ax.plot(sample_df[x_col], sample_df[y_col], color="#2563eb", linewidth=2)
        elif chart_type in ["Pie Distribution", "Pie Chart"]:
            pie_data = sample_df.groupby(x_col)[y_col].sum().head(8)
            ax.pie(pie_data, labels=pie_data.index, autopct='%1.1f%%', colors=sns.color_palette("pastel"))
        elif chart_type in ["Histogram Analytics", "Histogram"]:
            sns.histplot(data=sample_df, x=x_col, kde=True, color="#2563eb", ax=ax)
        elif chart_type in ["Box Plot Analysis", "Box Plot"]:
            sns.boxplot(data=sample_df, x=x_col, y=y_col, palette="Set2", ax=ax)
        elif chart_type in ["Heatmap Intensity", "Correlation Heatmap"]:
            numeric_df = sample_df.select_dtypes(include='number')
            if not numeric_df.empty and len(numeric_df.columns) > 1:
                sns.heatmap(numeric_df.corr(), annot=True, cmap="coolwarm", fmt=".2f", ax=ax)
            else:
                ax.text(0.5, 0.5, "Not Enough Numeric Data for Heatmap", ha='center', va='center')
        else:
            sns.barplot(data=sample_df, x=x_col, y=y_col, palette="Blues_d", ax=ax)

        plt.xticks(rotation=45, ha='right')
        plt.title(f"{chart_type} - Executive Export Summary", fontsize=14, fontweight='bold', pad=15)
        plt.tight_layout()
        
        img_buf = io.BytesIO()
        plt.savefig(img_buf, format='png', dpi=150)
        img_buf.seek(0)
        plt.close()
        return img_buf.getvalue()
    except Exception as e:
        return None


# ==========================================
# 3. PDF & MULTI-SLIDE PPTX REPORT ENGINES
# ==========================================
def create_pdf_report(df, include_graphs=False, graph_bytes=None):
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=16, style='B')
        pdf.cell(200, 10, txt="Arto+ Studio Cleaned Master Data Audit", ln=True, align='C')
        pdf.ln(10)
        
        pdf.set_font("Helvetica", size=11, style='B')
        pdf.cell(200, 10, txt=f"Total Records: {len(df)} | Attributes: {len(df.columns)}", ln=True)
        pdf.ln(5)
        
        pdf.set_font("Helvetica", size=8, style='B')
        cols_to_print = df.columns[:7]
        col_w = 190 / max(len(cols_to_print), 1)
        for col in cols_to_print: 
            pdf.cell(col_w, 10, str(col)[:12], border=1, align='C')
        pdf.ln()
        
        pdf.set_font("Helvetica", size=8)
        clean_df = df.fillna("-")
        for _, row in clean_df.iterrows():
            if pdf.get_y() > 270:
                pdf.add_page()
                pdf.set_font("Helvetica", size=8, style='B')
                for col in cols_to_print:
                    pdf.cell(col_w, 10, str(col)[:12], border=1, align='C')
                pdf.ln()
                pdf.set_font("Helvetica", size=8)
                
            for val in row[:7]:
                pdf.cell(col_w, 8, str(val)[:12], border=1, align='C')
            pdf.ln()
            
        if include_graphs and graph_bytes:
            pdf.add_page()
            pdf.set_font("Helvetica", size=14, style='B')
            pdf.cell(200, 10, txt="Data Visualization Analytics Chart", ln=True, align='C')
            pdf.ln(5)
            
            temp_pdf_img = "temp_report_chart.png"
            with open(temp_pdf_img, "wb") as f:
                f.write(graph_bytes)
            pdf.image(temp_pdf_img, x=15, y=30, w=180)
                
        return bytes(pdf.output(dest='S'))
    except Exception as e:
        return f"PDF Error: {str(e)}".encode('utf-8')


def create_ppt_report(df, include_graphs=False, graph_bytes=None):
    try:
        prs = Presentation()
        
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        slide1.shapes.title.text = "Arto+ Cleaned Data Presentation"
        txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        txBox.text_frame.text = f"Dataset Executive Audit Complete:\n\n• Cleaned Active Rows: {len(df)}\n• Columns Saved: {len(df.columns)}"
        
        rows_per_slide = 15
        num_cols = min(6, len(df.columns))
        clean_df = df.fillna("-")
        
        for start_idx in range(0, len(clean_df), rows_per_slide):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            title_box.text_frame.text = f"Master Cleaned Data Table (Rows {start_idx+1} to {min(start_idx+rows_per_slide, len(clean_df))})"
            
            chunk = clean_df.iloc[start_idx : start_idx + rows_per_slide]
            grid_rows = len(chunk) + 1
            
            table_shape = slide.shapes.add_table(grid_rows, num_cols, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.35 * grid_rows))
            table = table_shape.table
            
            for c_idx in range(num_cols):
                cell = table.cell(0, c_idx)
                cell.text = str(clean_df.columns[c_idx])[:15]
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.bold = True
                
            for r_idx, row in enumerate(chunk.values):
                for c_idx, val in enumerate(row[:num_cols]):
                    cell = table.cell(r_idx + 1, c_idx)
                    cell.text = str(val)[:20]
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        if include_graphs and graph_bytes:
            slide_g = prs.slides.add_slide(prs.slide_layouts[5])
            slide_g.shapes.title.text = "Automated Visual Analytics"
            
            temp_ppt_img = "temp_ppt_chart.png"
            with open(temp_ppt_img, "wb") as f:
                f.write(graph_bytes)
            slide_g.shapes.add_picture(temp_ppt_img, Inches(0.75), Inches(1.5), width=Inches(8.5))
                
        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()
    except Exception as e:
        return f"PPTX Error: {str(e)}".encode('utf-8')


# ==========================================
# 4. SIDEBAR GATEWAY WITH LICENCE INTERFACE
# ==========================================
if "master_df" not in st.session_state: st.session_state.master_df = None
if "original_raw_df" not in st.session_state: st.session_state.original_raw_df = None
if "selected_x" not in st.session_state: st.session_state.selected_x = None
if "selected_y" not in st.session_state: st.session_state.selected_y = None
if "selected_mode" not in st.session_state: st.session_state.selected_mode = "Bar Diagram"

# --- LICENCE KEY CHECK ---
VALID_KEYS = ["ARTO_PRO_2026", "VIP_ELITE_ACCESS"]
st.sidebar.markdown("### 👑 Elite Activation Suite")
user_key = st.sidebar.text_input("Enter Enterprise Licence Key:", type="password", help="Hint for testing: ARTO_PRO_2026")

is_premium = False
if user_key in VALID_KEYS:
    is_premium = True
    st.sidebar.success("👑 PRO UNLOCKED - Welcome Elite User!")
elif user_key:
    st.sidebar.error("❌ Invalid Activation Key Node.")

st.sidebar.markdown("---")
st.sidebar.markdown("### 📥 Data Management")
uploaded_files = st.sidebar.file_uploader("Upload Data Sheets (CSV, XLSX, JSON, TXT)", type=["csv", "xlsx", "xls", "json", "txt"], accept_multiple_files=True, key="sidebar_uploader")

if uploaded_files:
    imported_dfs = []
    for f in uploaded_files:
        try:
            if f.name.endswith('.csv') or f.name.endswith('.txt'): imported_dfs.append(pd.read_csv(f))
            elif f.name.endswith(('.xlsx', '.xls')): imported_dfs.append(pd.read_excel(f))
            elif f.name.endswith('.json'): imported_dfs.append(pd.read_json(f))
        except Exception as e:
            st.sidebar.error(f"Error reading: {f.name}")
            
    if imported_dfs:
        st.sidebar.markdown("#### 🔀 Structural Merging Layer")
        merge_axis = st.sidebar.radio("Merging Protocol Alignment:", ["Row-wise (Upar se Niche)", "Column-wise (Side-by-Side)"])
        
        if st.sidebar.button("⚡ Execute Pipeline Merge"):
            if merge_axis == "Row-wise (Upar se Niche)":
                merged_result = pd.concat(imported_dfs, axis=0, ignore_index=True)
            else:
                merged_result = pd.concat(imported_dfs, axis=1)
                
            st.session_state.master_df = merged_result.copy()
            st.session_state.original_raw_df = merged_result.copy()
            st.sidebar.success("⚡ Files Synchronized Panel Loaded!")
            st.rerun()

if st.session_state.master_df is not None:
    st.sidebar.markdown("---")
    if st.sidebar.button("🔄 System Hard Reset"):
        st.session_state.master_df = None
        st.session_state.original_raw_df = None
        st.session_state.selected_x = None
        st.session_state.selected_y = None
        st.rerun()

# TOP MAIN PREMIUM BRAND BAR WITH GLASS EFFECT LOOK
st.markdown("""
    <div class='brand-top-bar'>
        <div>
            <span class='brand-title'>Arto<sup>+</sup></span>
            <span class='brand-subtitle'>Automation Studio Pro</span>
        </div>
        <div class='badge-pro'>Enterprise Engine v3.0</div>
    </div>
""", unsafe_allow_html=True)


# ==========================================
# 5. WORKSPACE INTERACTIVE DESIGN TABS SYSTEM
# ==========================================
tabs = st.tabs([
    "✨ Data Workspace", 
    "🧼 Clean Data Core", 
    "🗑️ Duplicate Control", 
    "📈 Visual Analytics", 
    "🔍 Range Filtering", 
    "👁️ Final Audit View",   
    "📤 Enterprise Exports", 
    "📨 Dispatch Control"
])

if st.session_state.master_df is not None:
    working_df = st.session_state.master_df
    columns_list = list(working_df.columns)
    
    # ------------------------------------------
    # TAB 1: DATA WORKSPACE
    # ------------------------------------------
    with tabs[0]:
        st.markdown("### ✨ Interactive Live Sheet Fabricator")
        view_mode = st.radio("Data Matrix View Configuration:", ["Pura Data Dekhe (All Matrix)", "Row-wise Specific Filter", "Column-wise Specific Filter"], key="ws_view")
        display_df = working_df.copy()
        
        if view_mode == "Row-wise Specific Filter":
            selected_rows = st.multiselect("Select Row Index Subsets:", list(working_df.index))
            if selected_rows: display_df = working_df.loc[selected_rows]
        elif view_mode == "Column-wise Specific Filter":
            selected_cols = st.multiselect("Select Vector Attributes (Columns):", columns_list, default=columns_list[:5] if len(columns_list)>5 else columns_list)
            if selected_cols: display_df = working_df[selected_cols]
                
        updated_grid = st.data_editor(display_df, use_container_width=True, num_rows="dynamic", key="grid_editor")
        if not updated_grid.equals(display_df):
            if view_mode == "Pura Data Dekhe (All Matrix)": st.session_state.master_df = updated_grid
            elif view_mode == "Row-wise Specific Filter" and selected_rows: st.session_state.master_df.loc[selected_rows] = updated_grid
            elif view_mode == "Column-wise Specific Filter" and selected_cols: st.session_state.master_df[selected_cols] = updated_grid
            st.rerun()

    # ------------------------------------------
    # TAB 2: CLEAN DATA
    # ------------------------------------------
    with tabs[1]:
        st.markdown("### 🧼 AI-Assisted Cleansing Architecture")
        
        if st.session_state.original_raw_df is not None:
            st.markdown("#### ↩️ Smart History Pipeline Rollback")
            if st.button("⏪ Undo Changes & Bring Back Original None Values"):
                st.session_state.master_df = st.session_state.original_raw_df.copy()
                st.success("🔄 Rollback Complete! Original state with all None nodes has been re-established.")
                st.rerun()
            st.write("---")

        total_nulls = working_df.isna().sum().sum()
        st.metric(label="Detected Broken Nodes (None/Null Count)", value=total_nulls)
        
        if total_nulls > 0:
            st.markdown("⚡ **Absolute Structural Mapping of Missing Matrix Items:**")
            null_rows, null_cols = working_df.isna().values.nonzero()
            location_records = []
            for r, c in zip(null_rows, null_cols):
                location_records.append({
                    "Row coordinate (Index)": r, 
                    "Column Vector (Feature Name)": working_df.columns[c]
                })
            st.dataframe(pd.DataFrame(location_records), use_container_width=True)
            st.write("---")

        tc1, tc2 = st.columns(2)
        with tc1:
            case_col = st.selectbox("Text Transformation Column Selector:", ["None"] + columns_list)
            case_style = st.selectbox("Target Typography Format Style:", ["UPPERCASE", "lowercase", "Title Case"])
            if st.button("Execute Typography Conversion") and case_col != "None":
                if case_style == "UPPERCASE": st.session_state.master_df[case_col] = working_df[case_col].astype(str).str.upper()
                elif case_style == "lowercase": st.session_state.master_df[case_col] = working_df[case_col].astype(str).str.lower()
                else: st.session_state.master_df[case_col] = working_df[case_col].astype(str).str.title()
                st.rerun()
        with tc2:
            sort_target = st.selectbox("Structural Index Alignment Pivot Column:", ["None"] + columns_list, key="sort_c")
            sort_direction = st.radio("Sequence Array Target Direction:", ["Ascending (A-Z)", "Descending (Z-A)"])
            if st.button("Apply Spatial Sort Realignment") and sort_target != "None":
                st.session_state.master_df = working_df.sort_values(by=sort_target, ascending=(sort_direction == "Ascending (A-Z)")).reset_index(drop=True)
                st.rerun()
                
        st.write("---")
        
        st.markdown("#### 🛠️ Bulk Imputation Vector Management")
        nc1, nc2 = st.columns(2)
        with nc1:
            patch_val = st.text_input("Global Replacement Imputation Injection Value:")
            if st.button("Patch Across Entire Matrix Rows"):
                st.session_state.master_df = working_df.fillna(patch_val)
                st.rerun()
            if st.button("🚀 Rapid Multi-Impute All Blank Blocks with Integer Zero '0'"):
                st.session_state.master_df = working_df.fillna(0)
                st.rerun()
        with nc2:
            target_null_col = st.selectbox("Local Isolated Target Vector Channel:", ["None"] + columns_list)
            local_fill = st.text_input("Targeted Replace Entity Node String/Value:")
            if st.button("Apply Targeted Localized Imputation Patch") and target_null_col != "None":
                st.session_state.master_df[target_null_col] = working_df[target_null_col].fillna(local_fill)
                st.rerun()
                
        st.write("---")
        
        st.markdown("#### 🛑 Micro-Targeted Hard Structural Droppers (Row vs Column None Drop)")
        dc1, dc2 = st.columns(2)
        with dc1:
            st.warning("🚨 **Row Drop Action Engine**: Removes any record profile if a single element contains Null.")
            if st.button("🗑️ Purge Records / Rows Bearing Null Nodes"):
                st.session_state.master_df = working_df.dropna(axis=0).reset_index(drop=True)
                st.success("Targeted vector rows with missing items dropped from runtime environment.")
                st.rerun()
        with dc2:
            st.error("🚨 **Column Drop Action Engine**: Annihilates the entire attribute vector if any cells are missing.")
            if st.button("❌ Terminate Whole Columns Bearing Null Nodes"):
                st.session_state.master_df = working_df.dropna(axis=1)
                st.success("Entire attributes/columns with sparse missing nodes erased permanently.")
                st.rerun()

    # ------------------------------------------
    # TAB 3: DUPLICATE CONTROL
    # ------------------------------------------
    with tabs[2]:
        st.markdown("### 🗑️ Structural Mirror Integrity Firewall")
        dup_count = working_df.duplicated().sum()
        st.metric(label="Redundant Identical Rows Located", value=dup_count)
        
        if dup_count > 0:
            st.markdown("⚠️ **Deep Mirror Row Tracking Identification Interface:**")
            duplicate_indices = working_df[working_df.duplicated(keep=False)].index.tolist()
            st.write(f"Overlap clusters mapped inside these indexes: {duplicate_indices}")
            st.dataframe(working_df[working_df.duplicated(keep=False)], use_container_width=True)
            
            if st.button("Wipe Mirror Overlaps From Storage"):
                st.session_state.master_df = working_df.drop_duplicates().reset_index(drop=True)
                st.rerun()
        else:
            st.success("System Firewall Check Green: 100% unique row values detected.")

    # ------------------------------------------
    # TAB 4: VISUAL ANALYTICS
    # ------------------------------------------
    with tabs[3]:
        st.markdown("### 📈 Mathematical Vector Graphical Visualization")
        g_c1, g_c2 = st.columns(2)
        with g_c1:
            graph_scope = st.radio("Analytical Rendering Horizon Matrix Length:", ["Sabhi Data ka Graph Dekhe (All Records)", "Chuninda (Custom Row Range) Ka Graph Dekhe"])
        with g_c2:
            graph_col_filter = st.radio("Feature Dimensional Filtering Style:", ["Sabhi Columns Use Kare", "Sirf Selected Columns Matrix Use Kare"])
            
        chart_ready_df = working_df.copy()
        if graph_scope == "Chuninda (Custom Row Range) Ka Graph Dekhe":
            g_row_min, g_row_max = st.slider("Window Boundary Array Rows:", min_value=1, max_value=len(working_df), value=(1, min(15, len(working_df))))
            chart_ready_df = chart_ready_df.iloc[g_row_min-1:g_row_max]
            
        if graph_col_filter == "Sirf Selected Columns Matrix Use Kare":
            chosen_g_cols = st.multiselect("Chart Feature Parameters Input Vector Selection:", columns_list, default=columns_list[:2])
            if chosen_g_cols: chart_ready_df = chart_ready_df[chosen_g_cols]
                
        available_g_cols = list(chart_ready_df.columns)
        numeric_cols = chart_ready_df.select_dtypes(include='number').columns.tolist()
        
        vc1, vc2, vc3 = st.columns(3)
        with vc1: st.session_state.selected_x = st.selectbox("Primary Mapping Dimension (X-Axis):", available_g_cols)
        with vc2: st.session_state.selected_y = st.selectbox("Secondary Target Dependent Parameter (Y-Axis):", numeric_cols if numeric_cols else available_g_cols)
        with vc3:
            st.session_state.selected_mode = st.selectbox("Core Mathematical Visualization Element Module:", [
                "Bar Diagram", 
                "Line Plot Matrix", 
                "Scatter Cluster", 
                "Area Horizon", 
                "Pie Distribution", 
                "Histogram Analytics", 
                "Box Plot Analysis", 
                "Heatmap Intensity"
            ])
            
        try:
            cx, cy = st.session_state.selected_x, st.session_state.selected_y
            cmode = st.session_state.selected_mode
            
            if cmode == "Bar Diagram": fig = px.bar(chart_ready_df, x=cx, y=cy, color=cx, template="plotly_white", color_discrete_sequence=px.colors.qualitative.Prism)
            elif cmode == "Line Plot Matrix": fig = px.line(chart_ready_df, x=cx, y=cy, markers=True, template="plotly_white")
            elif cmode == "Scatter Cluster": fig = px.scatter(chart_ready_df, x=cx, y=cy, color=cx, size=cy if cy in numeric_cols else None, template="plotly_white", color_continuous_scale="Viridis")
            elif cmode == "Area Horizon": fig = px.area(chart_ready_df, x=cx, y=cy, template="plotly_white")
            elif cmode == "Pie Distribution": fig = px.pie(chart_ready_df, names=cx, values=cy if cy in numeric_cols else None, hole=0.4, template="plotly_white")
            elif cmode == "Histogram Analytics": fig = px.histogram(chart_ready_df, x=cx, kde=True, template="plotly_white")
            elif cmode == "Box Plot Analysis": fig = px.box(chart_ready_df, x=cx, y=cy, template="plotly_white")
            elif cmode == "Heatmap Intensity":
                if len(numeric_cols) > 1: fig = px.imshow(chart_ready_df[numeric_cols].corr(), text_auto=True, color_continuous_scale="Icefire")
                else: st.warning("More mathematical parameters needed to map correlation matrices."); fig = None
            
            if fig:
                fig.update_layout(margin=dict(l=20, r=20, t=40, b=20), paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
                st.plotly_chart(fig, use_container_width=True)
        except Exception as e:
            st.error(f"Visualization pipeline mapping exception: {e}")

    # ------------------------------------------
    # TAB 5: RANGE FILTERING
    # ------------------------------------------
    with tabs[4]:
        st.markdown("### 🔍 Mathematical Data Matrix Segment Slicer")
        row_min, row_max = st.slider("Window Index Array Selector Range:", min_value=1, max_value=len(working_df), value=(1, min(10, len(working_df))))
        filter_view_mode = st.radio("Segment Slicer Element Projection Options:", ["Pura Column Dekhe", "Chuninda Column Filter Kare"])
        
        sliced_df = working_df.iloc[row_min-1:row_max]
        if filter_view_mode == "Chuninda Column Filter Kare":
            range_cols = st.multiselect("Attribute Slicer Column Targets:", columns_list, default=columns_list)
            if range_cols: sliced_df = sliced_df[range_cols]
            
        st.dataframe(sliced_df, use_container_width=True)

    # ------------------------------------------
    # TAB 6: FINAL CLEANED PREVIEW
    # ------------------------------------------
    with tabs[5]:
        st.markdown("### 👁️ Structural Cleaned Master Data Audit Sheet")
        st.dataframe(working_df, use_container_width=True)

    # ------------------------------------------
    # TAB 7: REPORTS & EXPORTS (PREMIUM LOCK IMPLEMENTED)
    # ------------------------------------------
    with tabs[6]:
        st.markdown("### 📤 Enterprise Asset Compiling Hub")
        report_style = st.radio("Asset Payload Style Configuration Layout:", ["Graph ke sath (With Visual Analytics Chart)", "Without Graph (Only Clean Data Sheet)"])
        include_g = (report_style == "Graph ke sath (With Visual Analytics Chart)")
        
        ec1, ec2 = st.columns(2)
        with ec1:
            st.markdown("##### 🟢 Free Tier Formats")
            st.download_button("📥 Export Clean Data Array (.CSV)", data=working_df.to_csv(index=False).encode('utf-8'), file_name="Arto_Clean_Master.csv", mime="text/csv")
            
            excel_buffer = io.BytesIO()
            with pd.ExcelWriter(excel_buffer, engine='openpyxl') as wr:
                working_df.to_excel(wr, index=False)
            st.download_button("📥 Export Production Workbook (.XLSX)", data=excel_buffer.getvalue(), file_name="Arto_Clean_Master.xlsx")
            
            g_bytes = None
            if include_g and st.session_state.selected_x:
                g_bytes = generate_report_graph_bytes(working_df, st.session_state.selected_x, st.session_state.selected_y, st.session_state.selected_mode)

            pdf_binary = create_pdf_report(working_df, include_graphs=include_g, graph_bytes=g_bytes)
            st.download_button("📥 Compile Executive PDF Document", data=pdf_binary, file_name="Arto_Data_Report.pdf", mime="application/pdf")
            
        with ec2:
            st.markdown("##### 👑 Premium Corporate Asset")
            
            if not is_premium:
                st.markdown("""
                    <div class="premium-lock-container">
                        <div class="premium-lock-title">🔒 PPTX Presentation Locked</div>
                        <div class="premium-lock-desc">Upgrade to Pro Studio to auto-compile PowerPoint slide clusters.</div>
                    </div>
                """, unsafe_allow_html=True)
                
                # Locked UI Mock Element
                st.markdown('<div class="premium-blur-zone">', unsafe_allow_html=True)
                st.button("📥 Export Corporate Slide Deck (.PPTX)", key="locked_ppt_btn", disabled=True)
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                ppt_binary = create_ppt_report(working_df, include_graphs=include_g, graph_bytes=g_bytes)
                st.download_button("📥 Export Corporate Slide Deck (.PPTX)", data=ppt_binary, file_name="Arto_Presentation.pptx", key="unlocked_ppt_btn")

    # ------------------------------------------
    # TAB 8: BULK DISPATCHER (WHOLE TAB PREMIUM LOCKED)
    # ------------------------------------------
    with tabs[7]:
        st.markdown("### 📨 Omni-Channel Network Hub Message Transporter")
        
        if not is_premium:
            st.markdown("""
                <div class="premium-lock-container" style="max-width: 600px; margin: 40px auto;">
                    <div class="premium-lock-title" style="font-size: 22px;">🔒 Bulk Transmission Network Locked</div>
                    <p class="premium-lock-desc" style="font-size: 14px; margin-top: 10px;">
                        Automated Twilio SMS Gateways and SMTP Server Relays are reserved exclusively for premium subscribers.
                    </p>
                    <div style="font-size: 12px; color: #b45309; margin-top: 5px; font-weight:700;">
                        Activate via the side panel to unleash production workflows.
                    </div>
                </div>
            """, unsafe_allow_html=True)
            
            st.markdown('<div class="premium-blur-zone">', unsafe_allow_html=True)
            comm_mode = st.radio("Communication Protocol Channel:", ["SMTP Email Infrastructure", "Twilio SMS Network Gateway"], disabled=True)
            st.markdown('</div>', unsafe_allow_html=True)
        else:
            comm_mode = st.radio("Communication Protocol Channel Selector:", ["SMTP Email Infrastructure", "Twilio SMS Network Gateway"])
            
            if comm_mode == "SMTP Email Infrastructure":
                em1, em2 = st.columns(2)
                with em1:
                    src_email = st.text_input("SMTP Relay Authenticated Sender Address:")
                    src_pwd = st.text_input("SMTP Secure Access Hash Token:", type="password")
                with em2:
                    email_target_col = st.selectbox("Recipient Target Column Array Source:", columns_list)
                    mail_message_context = st.text_area("Mail Content Frame Plain Text Context:", "Greetings from Arto Automation Studio Hub!")
                    
                if st.button("Launch Network Broadcaster Mail Pipeline"):
                    if src_email and src_pwd:
                        try:
                            smtp_session = smtplib.SMTP("smtp.gmail.com", 587)
                            smtp_session.starttls()
                            smtp_session.login(src_email, src_pwd)
                            all_recipients = working_df[email_target_col].dropna().unique().tolist()
                            for recipient in all_recipients:
                                if "@" in str(recipient):
                                    packet = MIMEMultipart()
                                    packet['From'] = src_email
                                    packet['To'] = str(recipient)
                                    packet['Subject'] = "Enterprise Arto Studio Data Dispatch"
                                    packet.attach(MIMEText(mail_message_context, 'plain'))
                                    smtp_session.sendmail(src_email, str(recipient), packet.as_string())
                            smtp_session.quit()
                            st.success("Relay delivery handshake successful. All mail logs dispatched.")
                        except Exception as email_ex:
                            st.error(f"Transmission protocol failed: {email_ex}")
            else:
                sms1, sms2 = st.columns(2)
                with sms1:
                    t_sid = st.text_input("Twilio Endpoint Account Token SID Key:", type="password")
                    t_auth = st.text_input("Twilio Verification Security API Signature:", type="password")
                with sms2:
                    t_from = st.text_input("Registered Twilio Virtual Lease Outbound Number:")
                    phone_target_col = st.selectbox("Recipient Target Column Number Array Source:", columns_list)
                    sms_message_context = st.text_area("Telephony SMS Text Payload Frame:", "Arto Automation Server Notice Alert System Entry.")
                    
                if st.button("Initialize Gateway Broadcast SMS Outflow"):
                    if t_sid and t_auth:
                        try:
                            twilio_handler = Client(t_sid, t_auth)
                            destination_numbers = working_df[phone_target_col].dropna().unique().tolist()
                            for number in destination_numbers:
                                sanitized_num = str(number).strip()
                                if not sanitized_num.startswith('+'): sanitized_num = '+' + sanitized_num
                                twilio_handler.messages.create(body=sms_message_context, from_=t_from, to=sanitized_num)
                            st.success("Telephony payload injected. Gateway SMS stack cleared!")
                        except Exception as sms_ex:
                            st.error(f"Network node rejected SMS payload: {sms_ex}")
else:
    # WELCOME SCREEN CARD
    st.markdown("""
        <div style='text-align: center; margin-top: 60px; padding: 60px 40px; background: rgba(255,255,255,0.6); backdrop-filter: blur(12px); -webkit-backdrop-filter: blur(12px); border-radius: 24px; border: 2px dashed rgba(37,99,235,0.3); max-width: 800px; margin-left: auto; margin-right: auto; box-shadow: 0 20px 25px -5px rgba(0,0,0,0.02);'>
            <div style='background: linear-gradient(135deg, #eff6ff, #dbeafe); width: 80px; height: 80px; border-radius: 99px; display: flex; align-items: center; justify-content: center; margin: 0 auto 24px auto;'>
                <span style='font-size: 36px;'>📊</span>
            </div>
            <h3 style='font-size: 26px; color: #0f172a; font-weight:800; margin-bottom: 12px;'>Elite Enterprise Data Studio Ready</h3>
            <p style='color: #64748b; font-size: 16px; line-height: 1.6;'>Please ingest your file matrices via the left configuration sidebar deck. Choose your structural alignment protocol, and click <b>'Execute Pipeline Merge'</b> to spin up the UI workspace.</p>
        </div>
    """, unsafe_allow_html=True)